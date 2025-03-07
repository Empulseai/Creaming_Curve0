import pandas as pd
import matplotlib.pyplot as plt
import streamlit as st
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook
from openpyxl.styles import PatternFill,Font


st.set_page_config(page_title="Creaming Curve Analyzer")


st.markdown(
"""
    <style>
        body {
            background-color: #f4f4f4;
            font-family: 'Helvetica Neue', sans-serif;
        }
        .stApp {
            background-color: #f4f4f4;
        }
        h1 {
            color: #0078D7;
            text-align: center;
            font-size: 36px;
            font-family: 'Georgia', serif;
        }
        h4 {
            text-align: center;
            color: #444;
            font-family: 'Verdana', sans-serif;
        }
        .css-18e3th9 {
            padding-top: 2rem;
        }
        .stDataFrame {
            background-color: white;
            border-radius: 10px;
            padding: 10px;
        }
    </style>
    """,
    unsafe_allow_html=True
)


st.markdown(
    """
    <h1>Creaming Curve Analyzer</h1>
    <hr>
    """, 
    unsafe_allow_html=True
)


file=st.file_uploader("📂 Upload your Excel file", type=["xlsx"], help="Ensure the file contains relevant cost and savings data")


if file is not None:
    df = pd.read_excel(file)

    if 'Cost $' not in df.columns or 'Annual Savings $ K' not in df.columns:
        st.error("⚠️ Data is incomplete! Please enter 'Cost $' and 'Annual Savings $ K' values.")
        df['Cost $'] = 0
        df['Annual Savings $ K'] = 0  
    df = st.data_editor(df, num_rows="dynamic")

    df['Savings ratio'] = df['Annual Savings $ K'] / df['Cost $']
    df = df.sort_values(by='Savings ratio', ascending=False)
    df['Cumulative cost'] = df['Cost $'].cumsum()
    df['Cumulative Savings'] = df['Annual Savings $ K'].cumsum()


    st.markdown("Updated DataFrame:")
    st.dataframe(df.style.format({"Cost $": "${:,.2f}", "Annual Savings $ K": "${:,.2f}"}))

    fig, ax = plt.subplots(figsize=(12, 6))
    fig.patch.set_facecolor('#f4f4f4')

    budget = st.number_input("💰 Enter your budget ($ K):", min_value=0.0, step=100.0)


    if budget > 0:
        within_budget = df['Cumulative cost'] <= budget
        ax.scatter(df.loc[within_budget, 'Cumulative cost'], df.loc[within_budget, 'Cumulative Savings'], 
                   color='green', edgecolors='black', s=100, alpha=0.8, label='Projects within Budget')
        ax.scatter(df.loc[~within_budget, 'Cumulative cost'], df.loc[~within_budget, 'Cumulative Savings'], 
                   color='red', edgecolors='black', s=100, alpha=0.8, label='Projects outside Budget')
        ax.axvline(x=budget, color='red', linestyle='--', linewidth=2, label=f'Budget: ${budget:,.0f} K')
    else:
        ax.scatter(df['Cumulative cost'], df['Cumulative Savings'], 
                   color='blue', edgecolors='black', s=100, alpha=0.8, label='Projects')
         
    ax.set_title('Cumulative Cost vs. Cumulative Savings', fontsize=14, fontweight='bold', color='#333', fontname='Courier New')
    ax.set_xlabel('Cumulative Cost ($ K)', fontsize=10, fontweight='bold', color='#444', fontname='Tahoma')
    ax.set_ylabel('Cumulative Savings ($ K)', fontsize=10, fontweight='bold', color='#444', fontname='Tahoma')    

    ax.grid(True, linestyle='--', alpha=0.5, color='gray')

    x_labels = [f"{name}  (${cost:,.0f})" for name, cost in zip(df['Project Summary Name'], df['Cumulative cost'])]
    plt.xticks(df['Cumulative cost'], x_labels, rotation=90, ha='center', fontsize=8, color='#222', fontname='Arial')

    ax.legend(loc='upper left', frameon=True, facecolor='white', edgecolor='black')
    st.pyplot(fig)


    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title
    title.text = f"Creaming Curve                                 with Budget ${budget:,.0f}"
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(8))
    ppt_stream = BytesIO()
    prs.save(ppt_stream) 
    ppt_stream.seek(0)


    st.download_button("📥 Download PowerPoint", ppt_stream, "creaming_curve_presentation.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")

    excel_stream = BytesIO()
    with pd.ExcelWriter(excel_stream, engine='openpyxl') as writer:
     df.to_excel(writer, index=False, sheet_name='DataFrame')
     worksheet = writer.sheets['DataFrame']
    
     fill_green = PatternFill(start_color="90EE90", end_color="90EE90", fill_type="solid")  # Light green
     fill_red = PatternFill(start_color="FF9999", end_color="FF9999", fill_type="solid")  # Light red
     font_green = Font(color="006400", bold=False)  # White bold font for headers
     font_red = Font(color="8B0000", bold=False)  # White bold font for headers

    
     fill_blue = PatternFill(start_color="0000FF", end_color="0000FF", fill_type="solid")  # Blue header
     font_white_bold = Font(color="FFFFFF", bold=True)  # White bold font for headers

     for cell in worksheet[1]:  # First row (header row)
        cell.fill = fill_blue
        cell.font = font_white_bold
     
     for row in range(2, len(df) + 2):  
        is_within_budget = df.iloc[row - 2]['Cumulative cost'] <= budget
        fill_color = fill_green if is_within_budget else fill_red
        #fill_font=font_green if is_within_budget else font_red
        
        for col in worksheet.iter_cols(min_row=row, max_row=row, min_col=1, max_col=len(df.columns)):
            for cell in col:
                cell.fill = fill_color
                #cell.font=fill_font

    excel_stream.seek(0)
    
    st.download_button("📥 Download Excel", excel_stream, "creaming_curve_data.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    
    st.success("✅ Analysis Complete!")